"""Generate the presentation .pptx targeted at the rubric (PR1-PR4)."""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Inches, Pt

OUT = Path.home() / "Desktop" / "kickstarter_presentation_v2.pptx"

# colours
NAVY = RGBColor(0x12, 0x2A, 0x4F)
TEAL = RGBColor(0x0E, 0x9F, 0x6E)
ORANGE = RGBColor(0xD9, 0x77, 0x06)
RED = RGBColor(0xDC, 0x26, 0x26)
GREY = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0xF5, 0xF5, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    return s


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=NAVY, align="left"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = {"left": 1, "center": 2, "right": 3}[align]
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, *, size=18, color=NAVY):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = 1
        run = p.add_run()
        run.text = "•  " + item
        run.font.size = Pt(size)
        run.font.color.rgb = color
        p.space_after = Pt(8)
    return tb


def add_band(slide, color=NAVY, height=0.7):
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(height))
    band.line.fill.background()
    band.fill.solid()
    band.fill.fore_color.rgb = color


def add_footer(slide, n_total, n):
    add_text(
        slide, Inches(11.5), Inches(7.05), Inches(1.7), Inches(0.3),
        f"{n} / {n_total}", size=10, color=GREY, align="right",
    )
    add_text(
        slide, Inches(0.4), Inches(7.05), Inches(8), Inches(0.3),
        "Kickstarter Pre-Launch Recommender  ·  Omar Doughan",
        size=10, color=GREY,
    )


def title_slide():
    s = add_slide()
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.4), SW, Inches(2.7))
    band.line.fill.background()
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    add_text(
        s, Inches(0.6), Inches(2.6), Inches(12.2), Inches(1.2),
        "Kickstarter Pre-Launch Recommender",
        size=46, bold=True, color=WHITE,
    )
    add_text(
        s, Inches(0.6), Inches(3.7), Inches(12.2), Inches(1.0),
        "Pre-launch success prediction with SHAP-grounded creator advice",
        size=22, color=WHITE,
    )
    add_text(
        s, Inches(0.6), Inches(5.4), Inches(12.2), Inches(0.5),
        "Omar Doughan", size=18, bold=True, color=NAVY,
    )
    add_text(
        s, Inches(0.6), Inches(5.9), Inches(12.2), Inches(0.5),
        "XGBoost  ·  SHAP  ·  Streamlit  ·  Hugging Face Spaces",
        size=14, color=GREY,
    )


def header(slide, kicker, title):
    add_band(slide, NAVY, 0.7)
    add_text(slide, Inches(0.4), Inches(0.12), Inches(6), Inches(0.4),
             kicker, size=12, bold=True, color=WHITE)
    add_text(slide, Inches(0.4), Inches(0.85), Inches(12.5), Inches(0.7),
             title, size=30, bold=True, color=NAVY)


# ---- slide 1: title --------------------------------------------------------
title_slide()


# ---- slide 2: problem (PR1) ------------------------------------------------
s = add_slide()
header(s, "PR1  ·  PROBLEM", "Most Kickstarter campaigns fail. Creators want to know why.")
add_bullets(s, Inches(0.6), Inches(1.9), Inches(12), Inches(4.5), [
    "~40% of Kickstarter campaigns since 2009 failed to hit their funding goal.",
    "Existing 'success predictor' tools give a single number — no guidance.",
    "Goal: a tool a creator can use BEFORE launch to (1) estimate success probability "
    "and (2) get specific, feature-level advice on what to change.",
    "Constraint: predictions must use only what the creator controls at launch — "
    "no post-outcome features, no curatorial signals.",
], size=20)


# ---- slide 3: data (PR2) ---------------------------------------------------
s = add_slide()
header(s, "PR2  ·  DATA", "334k resolved campaigns, split by time")

# data box
def stat_card(slide, x, y, w, h, label, value, color=NAVY):
    bx = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    bx.line.color.rgb = color
    bx.line.width = Pt(1.5)
    bx.fill.solid()
    bx.fill.fore_color.rgb = LIGHT
    add_text(slide, x, y + Inches(0.15), w, Inches(0.5),
             value, size=28, bold=True, color=color, align="center")
    add_text(slide, x, y + Inches(0.85), w, Inches(0.4),
             label, size=12, color=GREY, align="center")

stat_card(s, Inches(0.6), Inches(2.0), Inches(2.9), Inches(1.4),
          "TOTAL CAMPAIGNS", "334,080")
stat_card(s, Inches(3.7), Inches(2.0), Inches(2.9), Inches(1.4),
          "TRAIN  (<2022)", "236,520", color=TEAL)
stat_card(s, Inches(6.8), Inches(2.0), Inches(2.9), Inches(1.4),
          "TEST  (≥2022)", "97,560", color=ORANGE)
stat_card(s, Inches(9.9), Inches(2.0), Inches(2.9), Inches(1.4),
          "FEATURE DIMS", "596,114")

add_text(s, Inches(0.6), Inches(3.7), Inches(12), Inches(0.5),
         "Source: WebRobots monthly Kickstarter scrapes (2009 – early 2026)",
         size=14, color=GREY)

add_bullets(s, Inches(0.6), Inches(4.4), Inches(12), Inches(2.8), [
    "Time-based split, not random — a random split leaks future info into training.",
    "Tabular block (52 dims): goal, duration, category, country, video flag, "
    "launch timing, blurb stats, Google Trends momentum.",
    "Text block (596,062 dims): word + character TF-IDF on the project blurb.",
], size=18)


# ---- slide 4: leakage (PR2) ------------------------------------------------
s = add_slide()
header(s, "PR2  ·  METHOD", "First model scored AUROC 1.00 — that was the bug, not the win")

box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(0.6), Inches(1.9), Inches(12.1), Inches(2.0))
box.line.color.rgb = RED
box.line.width = Pt(1.5)
box.fill.solid()
box.fill.fore_color.rgb = RGBColor(0xFE, 0xF2, 0xF2)
add_text(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(0.4),
         "LEAK FOUND", size=12, bold=True, color=RED)
add_text(s, Inches(0.9), Inches(2.4), Inches(11.5), Inches(1.5),
         "is_spotlight  →  Pearson correlation = 1.00 with the label.\n"
         "Kickstarter only sets that flag AFTER a campaign succeeds. Post-outcome.",
         size=18, color=NAVY)

add_text(s, Inches(0.6), Inches(4.2), Inches(12), Inches(0.5),
         "Two features dropped:", size=18, bold=True)
add_bullets(s, Inches(0.8), Inches(4.7), Inches(12), Inches(2.0), [
    "is_spotlight — strictly post-outcome.",
    "is_staff_pick — curatorial signal, not creator-controllable.",
    "Final feature set: 52 tabular + 596k TF-IDF, all controllable at launch.",
], size=18)


# ---- slide 5: models (PR2 / PR3) -------------------------------------------
s = add_slide()
header(s, "PR2 / PR3  ·  MODELS", "Baseline → tabular XGBoost → full text+tabular XGBoost")

# results table
rows = [
    ("Model", "Features", "Test AUROC"),
    ("Logistic Regression", "52 tabular", "0.773"),
    ("XGBoost (tabular)", "52", "0.826"),
    ("XGBoost (full)", "596,114", "0.868"),
]
tbl_left = Inches(0.6)
tbl_top = Inches(2.0)
tbl_w = Inches(12.1)
row_h = Inches(0.7)
col_widths = [Inches(5.0), Inches(4.6), Inches(2.5)]

x = tbl_left
for i, row in enumerate(rows):
    y = tbl_top + row_h * i
    is_head = i == 0
    is_winner = i == len(rows) - 1
    # row background
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, tbl_left, y, tbl_w, row_h)
    bg.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    bg.line.width = Pt(0.5)
    bg.fill.solid()
    if is_head:
        bg.fill.fore_color.rgb = NAVY
    elif is_winner:
        bg.fill.fore_color.rgb = RGBColor(0xDC, 0xFC, 0xE7)
    else:
        bg.fill.fore_color.rgb = WHITE
    cx = tbl_left
    for j, cell in enumerate(row):
        add_text(
            s, cx + Inches(0.2), y + Inches(0.15), col_widths[j] - Inches(0.4), Inches(0.5),
            cell, size=18,
            bold=is_head or is_winner or j == 2,
            color=WHITE if is_head else (TEAL if is_winner else NAVY),
        )
        cx += col_widths[j]

add_bullets(s, Inches(0.6), Inches(5.1), Inches(12), Inches(2.0), [
    "Tabular-only XGBoost = the deployable demo (8.5 MB model).",
    "Full model (596k features) trained on RTX 5070 with QuantileDMatrix to fit 12 GB VRAM.",
    "TF-IDF adds ~4 AUROC points — text helps, but tabular carries the bulk of the signal.",
], size=18)


# ---- slide 6: SHAP recommendations (PR2 / PR3) -----------------------------
s = add_slide()
header(s, "PR3  ·  CONTRIBUTION 1", "From a probability score to actionable advice")

add_text(s, Inches(0.6), Inches(1.9), Inches(12), Inches(0.5),
         "Translate per-feature SHAP attributions into plain-language advice.",
         size=20, color=NAVY)

# code-ish block
cx = Inches(0.6); cy = Inches(2.6); cw = Inches(12.1); ch = Inches(2.0)
codebox = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, cx, cy, cw, ch)
codebox.line.color.rgb = GREY
codebox.fill.solid()
codebox.fill.fore_color.rgb = LIGHT
add_text(s, cx + Inches(0.2), cy + Inches(0.15), cw - Inches(0.4), ch - Inches(0.3),
         "Rule fires only when:\n"
         "    SHAP[feature] < -0.05    AND    suggestion is valid for the current value\n"
         "Example: won't suggest 'add a video' if has_video already = 1.",
         size=16, color=NAVY)

add_text(s, Inches(0.6), Inches(4.8), Inches(12), Inches(0.5),
         "Sample output:", size=18, bold=True)

quote = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                          Inches(0.8), Inches(5.3), Inches(12), Inches(1.6))
quote.line.color.rgb = TEAL
quote.line.width = Pt(2)
quote.fill.solid()
quote.fill.fore_color.rgb = WHITE
add_text(s, Inches(1.0), Inches(5.4), Inches(11.6), Inches(1.5),
         "“Your goal of $250,000 may be too high for this profile. "
         "Consider lowering it.”\n"
         "“A 60-day campaign hurts your odds — most successful campaigns run 21–35 days.”",
         size=15, color=NAVY)


# ---- slide 7: temporal drift (PR3) -----------------------------------------
s = add_slide()
header(s, "PR3  ·  CONTRIBUTION 2", "A 2010-trained model is unsafe on 2022 campaigns")

# table of windows
rows = [
    ("Train window", "AUROC on own holdout", "AUROC on 2022+", "Δ"),
    ("2010 – 2014", "0.80", "0.64", "+0.16"),
    ("2015 – 2017", "0.81", "0.76", "+0.05"),
    ("2018 – 2020", "0.83", "0.79", "+0.04"),
    ("2021+",        "0.83", "0.86", "−0.03"),
]
tbl_left = Inches(0.6); tbl_top = Inches(1.9)
tbl_w = Inches(12.1); row_h = Inches(0.55)
col_widths = [Inches(3.0), Inches(3.6), Inches(3.0), Inches(2.5)]
for i, row in enumerate(rows):
    y = tbl_top + row_h * i
    is_head = i == 0
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, tbl_left, y, tbl_w, row_h)
    bg.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    bg.line.width = Pt(0.5)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY if is_head else WHITE
    cx = tbl_left
    for j, cell in enumerate(row):
        col = WHITE if is_head else NAVY
        if not is_head and j == 3:
            col = RED if cell.startswith("+") and float(cell[1:]) >= 0.05 else NAVY
            if cell.startswith("−"):
                col = TEAL
        add_text(
            s, cx + Inches(0.2), y + Inches(0.08), col_widths[j] - Inches(0.4), Inches(0.5),
            cell, size=16, bold=is_head, color=col,
        )
        cx += col_widths[j]

add_text(s, Inches(0.6), Inches(5.4), Inches(12), Inches(0.5),
         "Take-away:", size=18, bold=True)
add_bullets(s, Inches(0.8), Inches(5.85), Inches(12), Inches(1.5), [
    "Old data does not transfer — Kickstarter's distribution has shifted materially.",
    "We retrain on rolling windows. We do NOT trust a 2014 model on 2025 campaigns.",
], size=16)


# ---- slide 8: trends ablation (PR4 — owning the negative) ------------------
s = add_slide()
header(s, "PR4  ·  HONEST FINDING", "Google Trends momentum did not improve the model")

add_text(s, Inches(0.6), Inches(1.9), Inches(12), Inches(1.0),
         "Hypothesis: campaigns launched during topic spikes (per Google Trends) "
         "would outperform.", size=18, color=NAVY)

# ablation card
cy = Inches(3.0); ch = Inches(2.5)
box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(0.6), cy, Inches(12.1), ch)
box.line.color.rgb = ORANGE
box.line.width = Pt(1.5)
box.fill.solid()
box.fill.fore_color.rgb = LIGHT
add_text(s, Inches(0.9), cy + Inches(0.2), Inches(11.5), Inches(0.5),
         "ABLATION RESULTS", size=12, bold=True, color=ORANGE)
add_text(s, Inches(0.9), cy + Inches(0.7), Inches(11.5), Inches(0.5),
         "Full model with Google Trends      :   AUROC 0.868",
         size=20, color=NAVY)
add_text(s, Inches(0.9), cy + Inches(1.2), Inches(11.5), Inches(0.5),
         "Full model WITHOUT Google Trends :   AUROC 0.869",
         size=20, color=NAVY)
add_text(s, Inches(0.9), cy + Inches(1.8), Inches(11.5), Inches(0.5),
         "Difference: ≤ noise. Trends adds nothing on test.",
         size=18, bold=True, color=ORANGE)

add_text(s, Inches(0.6), Inches(5.9), Inches(12), Inches(1.0),
         "Reporting this honestly. The signal looked plausible a priori; "
         "the data did not support it.",
         size=15, color=GREY)


# ---- slide 9: demo (PR3) ---------------------------------------------------
s = add_slide()
header(s, "PR3  ·  DEMO", "Live on Hugging Face Spaces")

add_bullets(s, Inches(0.6), Inches(1.9), Inches(12), Inches(3.5), [
    "Streamlit app, Docker SDK, free CPU tier, private (collaborator access).",
    "Two modes:  (1) enter a hypothetical campaign,  (2) load a real held-out "
    "post-2022 campaign and predict its outcome.",
    "Output: probability  +  top SHAP contributors  +  natural-language advice.",
    "Ships the tabular XGBoost (8.5 MB).  Full 596k-feature model is too large "
    "for a free demo — we trade ~4 AUROC points for portability.",
], size=18)

url_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                            Inches(0.6), Inches(5.7), Inches(12.1), Inches(1.0))
url_box.line.color.rgb = TEAL
url_box.line.width = Pt(1.5)
url_box.fill.solid()
url_box.fill.fore_color.rgb = WHITE
add_text(s, Inches(0.8), Inches(5.85), Inches(11.7), Inches(0.7),
         "huggingface.co/spaces/OD2004/OD",
         size=24, bold=True, color=TEAL, align="center")


# ---- slide 9b: Lebanon case study (PR3 / PR4) ------------------------------
s = add_slide()
header(s, "PR3 / PR4  ·  CASE STUDY", "Lebanon-related campaigns — qualitative supplement")

add_text(s, Inches(0.6), Inches(1.85), Inches(12.1), Inches(0.8),
         "N = 49 is too small for reliable subgroup metrics "
         "(95% CI on AUROC ≈ 0.66–0.93). We go qualitative.",
         size=16, color=GREY)

# three cards
def case_card(slide, x, y, w, h, kicker, kicker_color, headline, body):
    bx = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    bx.line.color.rgb = kicker_color
    bx.line.width = Pt(1.5)
    bx.fill.solid()
    bx.fill.fore_color.rgb = WHITE
    add_text(slide, x + Inches(0.2), y + Inches(0.15), w - Inches(0.4), Inches(0.3),
             kicker, size=10, bold=True, color=kicker_color)
    add_text(slide, x + Inches(0.2), y + Inches(0.5), w - Inches(0.4), Inches(0.55),
             headline, size=15, bold=True, color=NAVY)
    add_text(slide, x + Inches(0.2), y + Inches(1.05), w - Inches(0.4), h - Inches(1.2),
             body, size=12, color=NAVY)

card_w = Inches(4.0); card_h = Inches(2.7); card_y = Inches(2.8)
case_card(s, Inches(0.6), card_y, card_w, card_h,
          "CORRECT RISK FLAG", RED,
          "Lebanese cuisine cookbook (AU)",
          "Goal $50k · no video · 7-word blurb.\n"
          "Model: 2% probability — flagged ALL three issues.\n"
          "Actual: failed. Tool would have given the creator "
          "concrete fixes pre-launch.")
case_card(s, Inches(4.7), card_y, card_w, card_h,
          "HIGH-CONFIDENCE WIN", TEAL,
          "AUTOMNE AU LEVANT (FR)",
          "Goal $3k · short campaign · video.\n"
          "Model: 96% probability.\n"
          "Actual: succeeded. Clean profile — model "
          "rewarded campaigns that got the basics right.")
case_card(s, Inches(8.8), card_y, card_w, card_h,
          "HUMILITY CHECK", ORANGE,
          "Surprise win",
          "Some campaigns succeed despite a low score — "
          "community trust and pre-launch outreach live "
          "outside the feature set. A low score is a "
          "signal to investigate, not a verdict.")

add_text(s, Inches(0.6), Inches(5.65), Inches(12.1), Inches(0.5),
         "Plus 3 hypothetical personas (relief film · cookbook · LB-based EP) — "
         "see lebanon_deep_case_study.md for the full walk-throughs.",
         size=14, color=GREY)


# ---- slide 10: limitations (PR4) -------------------------------------------
s = add_slide()
header(s, "PR4  ·  LIMITATIONS", "What this tool is NOT")

add_bullets(s, Inches(0.6), Inches(1.9), Inches(12.1), Inches(5.0), [
    "Correlational, not causal. SHAP says which features the model used; it does "
    "NOT say changing them would change the outcome.",
    "Per-group AUROC ranges 0.69 – 0.90 across categories and countries. "
    "The model is not equally calibrated everywhere.",
    "Base-rate shift: post-2022 campaigns succeed 72% in this scrape vs 61% pre-2022. "
    "Could be platform changes; could be scraping selection bias. We don't claim to know.",
    "Deployed model is the tabular XGBoost (0.826 AUROC). The headline 0.868 number "
    "uses 596k features and is not in the demo.",
    "Threshold tuning: balanced operating point at p = 0.64 (P = 0.85, R = 0.86). "
    "A different threshold is appropriate if false-positive cost ≠ false-negative cost.",
], size=17)


# ---- slide 11: Q&A / closer (PR4) -----------------------------------------
s = add_slide()
header(s, "PR4  ·  WRAP", "What I'd do next")

add_bullets(s, Inches(0.6), Inches(1.9), Inches(12), Inches(4.0), [
    "Calibrate per-category. The 0.69–0.90 spread across categories means a single "
    "global threshold is the wrong call for some segments.",
    "Re-test the Trends signal with a query-specific lookup, not a generic one — "
    "the current ablation may have under-powered it.",
    "Causal layer: replace SHAP-based 'recommendations' with an uplift model trained "
    "on near-duplicate campaigns that differ only in one feature.",
    "Productionise the retraining cadence — temporal drift implies a 12-month rebuild "
    "schedule at minimum.",
], size=18)

# closer
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                         Inches(0.6), Inches(6.0), Inches(12.1), Inches(0.9))
band.line.fill.background()
band.fill.solid()
band.fill.fore_color.rgb = NAVY
add_text(s, Inches(0.6), Inches(6.15), Inches(12.1), Inches(0.6),
         "Questions?",
         size=28, bold=True, color=WHITE, align="center")


# footers + save
n_total = len(prs.slides)
for i, slide in enumerate(prs.slides):
    if i == 0:
        continue
    add_footer(slide, n_total, i + 1)

prs.save(OUT)
print(f"saved: {OUT}")
print(f"slides: {n_total}")
